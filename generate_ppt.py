"""
路演PPT生成器 — iOS 26 Apple 风格
基于多源工业数据融合的智能预测性维护决策平台
输出：7页 .pptx 文件，浅色玻璃拟态 + 环境光晕 + 大圆角
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
import random

# ══ iOS 26 Design Tokens (Light Mode) ══
BG_ROOT     = RGBColor(0xF2, 0xF2, 0xF7)
BG_SURFACE  = RGBColor(0xFA, 0xFA, 0xFA)
TEXT_PRI    = RGBColor(0x1D, 0x1D, 0x1F)  # primary
TEXT_SEC    = RGBColor(0x6E, 0x6E, 0x73)  # secondary
TEXT_MUTED  = RGBColor(0xAE, 0xAE, 0xB2)  # muted

# Apple semantic palette
CYAN    = RGBColor(0x5A, 0xC8, 0xB8)
GREEN   = RGBColor(0x4C, 0xD9, 0x64)
BLUE    = RGBColor(0x64, 0xB5, 0xF6)
AMBER   = RGBColor(0xF0, 0xA8, 0x40)
RED     = RGBColor(0xE0, 0x60, 0x60)
PURPLE  = RGBColor(0xB3, 0x88, 0xEB)
PINK    = RGBColor(0xE0, 0x80, 0xA0)

# Glass card (light mode): rgba(255,255,255,0.38) → ~RGB(253,253,253)
GLASS_BG      = RGBColor(0xFD, 0xFD, 0xFD)  # simulates ~38% white on f2f2f7
GLASS_BG_DEEP = RGBColor(0xFA, 0xFA, 0xFA)  # slightly more opaque
GLASS_BORDER  = RGBColor(0xE5, 0xE5, 0xEA)  # rgba(0,0,0,0.06) equiv
GLASS_BORDER_ACCENT = RGBColor(0xD1, 0xD1, 0xD6)  # rgba(0,0,0,0.10)
SHADOW_SM = RGBColor(0x00, 0x00, 0x00)  # shadow base, used lightly
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

FONT_SANS = 'PingFang SC'
FONT_BOLD = 'PingFang SC'
FONT_NUM  = 'SF Mono'

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# ═══════════════════════════════════════════
# HELPER: XML shadow effect on a shape
# ═══════════════════════════════════════════
def add_shadow(shape, blur_pt=8, dist_pt=2, alpha_pct=6000):
    """Add a subtle outer shadow to a shape via XML."""
    spPr = shape._element.find(qn('a:spPr'))
    if spPr is None:
        spPr = shape._element.makeelement(qn('a:spPr'), {})
        shape._element.insert(0, spPr)
    effectLst = spPr.find(qn('a:effectLst'))
    if effectLst is None:
        effectLst = spPr.makeelement(qn('a:effectLst'), {})
        spPr.append(effectLst)
    outerShdw = effectLst.makeelement(qn('a:outerShdw'), {
        'blurRad': str(int(blur_pt * 12700)),
        'dist': str(int(dist_pt * 12700)),
        'dir': '5400000',
        'algn': 'ctr',
    })
    srgbClr = outerShdw.makeelement(qn('a:srgbClr'), {'val': '000000'})
    alpha = srgbClr.makeelement(qn('a:alpha'), {'val': str(alpha_pct)})
    srgbClr.append(alpha)
    outerShdw.append(srgbClr)
    effectLst.append(outerShdw)

# ═══════════════════════════════════════════
# Ambient orb: simulate radial gradient glow
# ═══════════════════════════════════════════
def add_ambient_orb(slide, left, top, width, height, color, alpha_val=6):
    """Large semi-transparent ellipse for background ambient light."""
    orb = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
    orb.fill.solid()
    rgb_str = '{:02X}{:02X}{:02X}'.format(color[0], color[1], color[2])
    orb.fill.fore_color.rgb = RGBColor(color[0], color[1], color[2])
    # Apply transparency via shape XML
    spPr = orb._element.find(qn('a:spPr'))
    if spPr is None:
        spPr = orb._element.makeelement(qn('a:spPr'), {})
        orb._element.insert(0, spPr)
    solidFill_el = spPr.find(qn('a:solidFill'))
    if solidFill_el is None:
        solidFill_el = spPr.makeelement(qn('a:solidFill'), {})
        spPr.append(solidFill_el)
    srgbClr = solidFill_el.find(qn('a:srgbClr'))
    if srgbClr is None:
        srgbClr = solidFill_el.makeelement(qn('a:srgbClr'), {'val': rgb_str})
        solidFill_el.append(srgbClr)
    else:
        srgbClr.set('val', rgb_str)
    # Remove existing alpha elements
    for a in srgbClr.findall(qn('a:alpha')):
        srgbClr.remove(a)
    alpha_el = srgbClr.makeelement(qn('a:alpha'), {'val': str(alpha_val * 1000)})
    srgbClr.append(alpha_el)
    orb.line.fill.background()
    return orb

# ═══════════════════════════════════════════
# Helper: set slide bg
# ═══════════════════════════════════════════
def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

# ═══════════════════════════════════════════
# Helper: add textbox
# ═══════════════════════════════════════════
def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=TEXT_PRI, bold=False, alignment=PP_ALIGN.LEFT,
                font_name=None, line_spacing=1.2):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name or FONT_SANS
    p.alignment = alignment
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    pPr = p._pPr
    if pPr is None:
        pPr = p._p.get_or_add_pPr()
    lnSpc = pPr.makeelement(qn('a:lnSpc'), {})
    spcPct = lnSpc.makeelement(qn('a:spcPct'), {'val': str(int(line_spacing * 100000))})
    lnSpc.append(spcPct)
    pPr.append(lnSpc)
    return txBox

def add_multiline_textbox(slide, left, top, width, height, lines,
                          font_size=18, color=TEXT_PRI, bold=False,
                          alignment=PP_ALIGN.LEFT, font_name=None, line_spacing=1.5):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(lines):
        if isinstance(item, str):
            text, overrides = item, {}
        else:
            text, overrides = item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        fs = overrides.get('font_size', font_size)
        fc = overrides.get('color', color)
        fb = overrides.get('bold', bold)
        fn = overrides.get('font_name', font_name or FONT_SANS)
        fa = overrides.get('alignment', alignment)
        p.font.size = Pt(fs)
        p.font.color.rgb = fc
        p.font.bold = fb
        p.font.name = fn
        p.alignment = fa
        p.space_after = Pt(overrides.get('space_after', 2))
        p.space_before = Pt(overrides.get('space_before', 0))
        ls_val = overrides.get('line_spacing', line_spacing)
        pPr = p._pPr
        if pPr is None:
            pPr = p._p.get_or_add_pPr()
        lnSpc = pPr.makeelement(qn('a:lnSpc'), {})
        spcPct = lnSpc.makeelement(qn('a:spcPct'), {'val': str(int(ls_val * 100000))})
        lnSpc.append(spcPct)
        pPr.append(lnSpc)
    return txBox

# ═══════════════════════════════════════════
# Glass Card
# ═══════════════════════════════════════════
def add_glass_card(slide, left, top, width, height, corner_radius=Inches(0.15), shadow=True, fill_override=None):
    """iOS 26 glass card: white-semi-transparent + subtle border + shadow + large radius."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_override if fill_override else GLASS_BG
    shape.line.color.rgb = GLASS_BORDER
    shape.line.width = Pt(0.5)
    # Large corner radius via XML
    spPr = shape._element.find(qn('a:spPr'))
    if spPr is not None:
        prstGeom = spPr.find(qn('a:prstGeom'))
        if prstGeom is not None:
            prstGeom.set('prst', 'roundRect')
        # Set corner radius via avLst
        for avLst in prstGeom.findall(qn('a:avLst')):
            prstGeom.remove(avLst)
    if shadow:
        add_shadow(shape, blur_pt=10, dist_pt=1, alpha_pct=4000)
    return shape

def _rgb_tuple(c):
    """Extract (r,g,b) from python-pptx RGBColor or hex string."""
    s = str(c)
    if s.startswith('#'):
        s = s[1:]
    if len(s) == 6:
        return (int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))
    return (0xD0, 0xD0, 0xD0)  # fallback

def add_glass_card_colored(slide, left, top, width, height, accent_color, corner_radius=Inches(0.15)):
    """Glass card with a very subtle colored tint from the accent color."""
    ar, ag, ab = _rgb_tuple(accent_color)
    # Mix: 6% accent + 94% white
    tr = min(255, int(ar * 0.06 + 253))
    tg = min(255, int(ag * 0.06 + 253))
    tb = min(255, int(ab * 0.06 + 253))
    tint = RGBColor(tr, tg, tb)

    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = tint
    shape.line.color.rgb = GLASS_BORDER
    shape.line.width = Pt(0.5)
    add_shadow(shape, blur_pt=8, dist_pt=1, alpha_pct=3000)
    return shape

# ═══════════════════════════════════════════
# Image placeholder (Apple style)
# ═══════════════════════════════════════════
def add_image_placeholder(slide, left, top, width, height, label="", sublabel=""):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xFA)
    shape.line.color.rgb = GLASS_BORDER
    shape.line.width = Pt(0.5)
    add_shadow(shape, blur_pt=6, dist_pt=1, alpha_pct=2000)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = label if label else ""
    p.font.size = Pt(13)
    p.font.color.rgb = TEXT_MUTED
    p.font.name = FONT_SANS
    p.alignment = PP_ALIGN.CENTER
    if sublabel:
        p2 = tf.add_paragraph()
        p2.text = sublabel
        p2.font.size = Pt(10)
        p2.font.color.rgb = TEXT_MUTED
        p2.font.name = FONT_SANS
        p2.alignment = PP_ALIGN.CENTER
    # Vertically center first para
    if label:
        p.space_before = Pt(max(0, int(height/Pt(1)/2) - 30))
    return shape

# ═══════════════════════════════════════════
# Page number
# ═══════════════════════════════════════════
def add_page_number(slide, num):
    add_textbox(slide, Inches(12.3), Inches(7.1), Inches(0.7), Inches(0.3),
                f'{num:02d}', font_size=10, color=TEXT_MUTED,
                alignment=PP_ALIGN.RIGHT, font_name=FONT_NUM)

# ═══════════════════════════════════════════
# Background orbs for all content slides
# ═══════════════════════════════════════════
def add_ambient_orbs(slide):
    # Mirror design-tokens.css body::before:
    # cyan orb top-left, blue orb bottom-right, purple orb bottom-center
    add_ambient_orb(slide, Inches(-1), Inches(-1), Inches(6), Inches(4), (90,200,184), alpha_val=8)
    add_ambient_orb(slide, Inches(8), Inches(4), Inches(6), Inches(4), (100,181,246), alpha_val=6)
    add_ambient_orb(slide, Inches(4), Inches(5), Inches(5), Inches(3), (179,136,235), alpha_val=5)


# ═══════════════════════════════════════════
# PAGE 1 — 封面 (Cover)
# ═══════════════════════════════════════════
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide1, BG_ROOT)

# Ambient orbs for depth
add_ambient_orb(slide1, Inches(-1.5), Inches(-1.5), Inches(8), Inches(5.5), (90,200,184), alpha_val=10)
add_ambient_orb(slide1, Inches(7), Inches(3.5), Inches(7), Inches(5), (100,181,246), alpha_val=7)
add_ambient_orb(slide1, Inches(3), Inches(5.5), Inches(6), Inches(3), (179,136,235), alpha_val=5)

# Hero glass card — centered
hero_card = add_glass_card(slide1, Inches(1.8), Inches(1.5), Inches(9.7), Inches(4.6), corner_radius=Inches(0.25))

# Tiny accent line above title
line_shape = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.7), Inches(2.2), Inches(2.0), Pt(2.5))
line_shape.fill.solid()
line_shape.fill.fore_color.rgb = CYAN
line_shape.line.fill.background()

# Main title
add_textbox(slide1, Inches(2.5), Inches(2.5), Inches(8.3), Inches(1.2),
            '基于多源工业数据融合的\n智能预测性维护决策平台',
            font_size=36, color=TEXT_PRI, bold=True, alignment=PP_ALIGN.CENTER, line_spacing=1.35)

# Subtitle
add_textbox(slide1, Inches(2.5), Inches(3.8), Inches(8.3), Inches(0.6),
            '从"数据瓶颈"到"决策闭环"',
            font_size=22, color=CYAN, bold=False, alignment=PP_ALIGN.CENTER)

# Separator
sep_shape = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.0), Inches(4.55), Inches(3.3), Pt(0.5))
sep_shape.fill.solid()
sep_shape.fill.fore_color.rgb = GLASS_BORDER
sep_shape.line.fill.background()

# Team & track
add_textbox(slide1, Inches(2.5), Inches(4.8), Inches(8.3), Inches(0.4),
            '××× 团队  ·  工业智能运维赛道',
            font_size=16, color=TEXT_SEC, alignment=PP_ALIGN.CENTER)

# Date
add_textbox(slide1, Inches(2.5), Inches(5.25), Inches(8.3), Inches(0.3),
            '2026.06',
            font_size=13, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)

# Logo placeholder
add_image_placeholder(slide1, Inches(10.8), Inches(6.4), Inches(1.8), Inches(0.7),
                      "Logo")

add_page_number(slide1, 1)


# ═══════════════════════════════════════════
# PAGE 2 — 背景与痛点
# ═══════════════════════════════════════════
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide2, BG_ROOT)
add_ambient_orbs(slide2)

# Page title (SF-style: bold, clean)
add_textbox(slide2, Inches(0.9), Inches(0.45), Inches(4), Inches(0.5),
            '背景与痛点', font_size=28, color=TEXT_PRI, bold=True)

# Headline
add_textbox(slide2, Inches(0.9), Inches(1.05), Inches(7), Inches(0.45),
            '100台CNC机床的被动维护模式，正在吃掉工厂的利润和产能。',
            font_size=19, color=AMBER, bold=True)

# Scene info card
add_glass_card(slide2, Inches(0.9), Inches(1.7), Inches(5.5), Inches(0.65), corner_radius=Inches(0.12))
add_textbox(slide2, Inches(1.2), Inches(1.78), Inches(5.0), Inches(0.45),
            '100 台 CNC 数控机床  ·  4 维传感器（电压 / 电流 / 温度 / 转速）  ·  9 种故障类型',
            font_size=14, color=TEXT_PRI, bold=True)

# Three pain point cards
pain_data = [
    ('01', '被动维修', '坏了才修\n非计划停机损失\n产能 15%–30%', RED),
    ('02', '人工巡检', '经验驱动\n不同班组标准不一\n漏检率难以控制', AMBER),
    ('03', '策略缺失', '成本 / 效率 / 质量\n无法量化权衡\n缺乏工具支撑决策', BLUE),
]
for idx, (num, title, desc, color) in enumerate(pain_data):
    left = Inches(0.9 + idx * 2.15)
    top = Inches(2.6)
    card = add_glass_card_colored(slide2, left, top, Inches(1.95), Inches(2.3), color)
    # Number
    add_textbox(slide2, left + Inches(0.18), top + Inches(0.15), Inches(0.5), Inches(0.35),
                num, font_size=26, color=color, bold=True, font_name=FONT_NUM)
    # Title
    add_textbox(slide2, left + Inches(0.18), top + Inches(0.65), Inches(1.6), Inches(0.35),
                title, font_size=16, color=TEXT_PRI, bold=True)
    # Description
    add_textbox(slide2, left + Inches(0.18), top + Inches(1.15), Inches(1.6), Inches(0.95),
                desc, font_size=12, color=TEXT_SEC, line_spacing=1.45)

# Core question
add_textbox(slide2, Inches(0.9), Inches(5.2), Inches(5.5), Inches(0.4),
            '核心命题', font_size=12, color=TEXT_MUTED, bold=True, font_name=FONT_NUM)

cq_card = add_glass_card_colored(slide2, Inches(0.9), Inches(5.55), Inches(5.5), Inches(0.65), GREEN)
add_textbox(slide2, Inches(1.2), Inches(5.68), Inches(5.0), Inches(0.4),
            '如何用数据驱动把"被动抢修" → "主动防御"？',
            font_size=17, color=TEXT_PRI, bold=True)

# Arrow emphasis
add_textbox(slide2, Inches(2.8), Inches(6.4), Inches(2.5), Inches(0.35),
            '被动抢修  ──→  主动防御',
            font_size=14, color=GREEN, alignment=PP_ALIGN.CENTER, font_name=FONT_NUM)

# Right: image placeholder
add_image_placeholder(slide2, Inches(7.3), Inches(1.7), Inches(5.2), Inches(4.9),
                      "CNC 机床产线实拍图",
                      "暗调处理，不抢文字焦点")

add_page_number(slide2, 2)


# ═══════════════════════════════════════════
# PAGE 3 — 关键发现：数据瓶颈
# ═══════════════════════════════════════════
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide3, BG_ROOT)
add_ambient_orbs(slide3)

add_textbox(slide3, Inches(0.9), Inches(0.45), Inches(6), Inches(0.5),
            '关键发现：数据瓶颈', font_size=28, color=TEXT_PRI, bold=True)

# Headline — this is the turning point page
add_textbox(slide3, Inches(0.9), Inches(1.05), Inches(11.5), Inches(0.45),
            '仅靠4个传感器做纯ML预测，信息量不够——这不是模型的问题，是输入的问题。',
            font_size=19, color=RED, bold=True)

# Three conclusion cards
conc_data = [
    ("Youden's J ≤ 0.075", '理论 AUC 上限 ≈ 0.537\n4 参数信息量接近随机猜测', RED),
    ('7 算法实测 AUC', 'RF / XGBoost / LightGBM\nSVM / LR / KNN / MLP\n全部落在 0.517 – 0.589', AMBER),
    ('设备间方差 61%–73%', '全局阈值完全无效\n每台设备必须建立\n独立统计基线', BLUE),
]
for idx, (title, desc, color) in enumerate(conc_data):
    left = Inches(0.9 + idx * 2.55)
    top = Inches(1.7)
    card = add_glass_card_colored(slide3, left, top, Inches(2.35), Inches(2.4), color)
    add_textbox(slide3, left + Inches(0.2), top + Inches(0.18), Inches(1.95), Inches(0.5),
                title, font_size=17, color=color, bold=True, font_name=FONT_NUM)
    add_textbox(slide3, left + Inches(0.2), top + Inches(0.85), Inches(1.95), Inches(1.3),
                desc, font_size=12, color=TEXT_SEC, line_spacing=1.5)

# Pivot callout
pivot_card = add_glass_card_colored(slide3, Inches(0.9), Inches(4.4), Inches(7.2), Inches(0.85), RED)
add_textbox(slide3, Inches(1.2), Inches(4.52), Inches(6.6), Inches(0.4),
            '我们的判断：不堆模型，换路线。', font_size=21, color=RED, bold=True)
add_textbox(slide3, Inches(1.2), Inches(4.92), Inches(6.6), Inches(0.3),
            '→ 转向 多信号融合 + 统计基线 方案', font_size=15, color=AMBER, font_name=FONT_NUM)

# Charts: left KDE, right AUC bars
add_image_placeholder(slide3, Inches(0.9), Inches(5.55), Inches(5.8), Inches(1.6),
                      "KDE 分布叠加图：4 参数 × 200 点分布曲线",
                      "突出低分离度  ·  标注 Youden's J ≤ 0.075")
add_image_placeholder(slide3, Inches(7.1), Inches(5.55), Inches(5.4), Inches(1.6),
                      "7 算法 AUC 对比柱状图",
                      "红色虚线标注随机基线 0.5  ·  AUC 0.517 – 0.589")

add_page_number(slide3, 3)


# ═══════════════════════════════════════════
# PAGE 4 — 解决方案架构
# ═══════════════════════════════════════════
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide4, BG_ROOT)
add_ambient_orbs(slide4)

add_textbox(slide4, Inches(0.9), Inches(0.45), Inches(6), Inches(0.5),
            '解决方案架构', font_size=28, color=TEXT_PRI, bold=True)

add_textbox(slide4, Inches(0.9), Inches(1.05), Inches(11.5), Inches(0.45),
            '统计基线发现异常，多信号融合做出决策——一套会解释"为什么"的工业决策引擎。',
            font_size=19, color=AMBER, bold=True)

# Central architecture diagram placeholder
add_image_placeholder(slide4, Inches(0.9), Inches(1.7), Inches(11.5), Inches(2.8),
                      "系统架构流程图  —  中央大幅",
                      "数据采集 → 逐设备统计基线（Z-Score + T²）→ 四路信号融合 → SHAP 解释 → 决策输出 → 四级降级 → 工单执行")

# Four layer cards
layer_data = [
    ('Layer 1', '逐设备统计基线', 'Z-Score + Hotelling T²\n≥ 3 样本建基线\n三层回退策略', CYAN),
    ('Layer 2', '多信号融合决策', '统计异常 + ML 概率\n+ 成本风险矩阵\n+ 趋势衰减预测', AMBER),
    ('Layer 3', 'SHAP 可解释性', 'RiskDecomposer →\nLocalExplainer\n告警 → 参数 → 根因', GREEN),
    ('Layer 4', '四级降级保障', 'FULL → STAT_ONLY\n→ RULE_ONLY\n→ EMERGENCY', PURPLE),
]
for idx, (layer, title, desc, color) in enumerate(layer_data):
    left = Inches(0.9 + idx * 3.05)
    top = Inches(4.8)
    card = add_glass_card_colored(slide4, left, top, Inches(2.8), Inches(2.3), color)
    add_textbox(slide4, left + Inches(0.2), top + Inches(0.12), Inches(1.0), Inches(0.25),
                layer, font_size=10, color=color, bold=True, font_name=FONT_NUM)
    add_textbox(slide4, left + Inches(0.2), top + Inches(0.4), Inches(2.4), Inches(0.35),
                title, font_size=15, color=TEXT_PRI, bold=True)
    add_textbox(slide4, left + Inches(0.2), top + Inches(0.85), Inches(2.4), Inches(1.2),
                desc, font_size=12, color=TEXT_SEC, line_spacing=1.5)

add_page_number(slide4, 4)


# ═══════════════════════════════════════════
# PAGE 5 — 平台工程实证
# ═══════════════════════════════════════════
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide5, BG_ROOT)
add_ambient_orbs(slide5)

add_textbox(slide5, Inches(0.9), Inches(0.45), Inches(6), Inches(0.5),
            '平台工程实证', font_size=28, color=TEXT_PRI, bold=True)

add_textbox(slide5, Inches(0.9), Inches(1.05), Inches(11.5), Inches(0.45),
            '从数据到决策，全链路闭环——12 个前端页面 + 三种角色 + 三种策略，一套平台全部跑通。',
            font_size=19, color=AMBER, bold=True)

# Left: main dashboard screenshot
add_image_placeholder(slide5, Inches(0.9), Inches(1.7), Inches(6.5), Inches(3.8),
                      "仪表盘主页截图  —  大图",
                      "工单总览 + 设备状态面板 + ECharts 可视化")

# Right top: 3-role comparison
add_image_placeholder(slide5, Inches(7.8), Inches(1.7), Inches(4.8), Inches(1.8),
                      "三角色视图对比：运维 / 管理 / 开发",
                      "同一页面，CSS 角色滤镜驱动不同视图 · 三张截图横排")

# Right bottom: work order tracking
add_image_placeholder(slide5, Inches(7.8), Inches(3.8), Inches(4.8), Inches(1.7),
                      "工单追踪页面截图",
                      "6 状态闭环：待处理 → 已派发 → 进行中 → 待验收 → 已完成 → 已归档")

# Feature badge row
features = ['12 前端页面', '3 角色分版', '3 维护策略', '6 状态工单', 'RAG 知识库', '故障注入', '40+ 话术库']
for idx, feat in enumerate(features):
    left = Inches(0.9 + idx * 1.75)
    badge = add_glass_card(slide5, left, Inches(5.85), Inches(1.6), Inches(0.45), corner_radius=Inches(0.08), shadow=False)
    add_textbox(slide5, left + Inches(0.05), Inches(5.88), Inches(1.5), Inches(0.4),
                feat, font_size=13, color=TEXT_PRI, bold=True, alignment=PP_ALIGN.CENTER)

# Bottom note
add_textbox(slide5, Inches(0.9), Inches(6.6), Inches(11.5), Inches(0.4),
            '以上均为已跑通的代码功能，非 PPT 概念展示。评委可现场体验故障注入演示与评委讲解助手。',
            font_size=12, color=GREEN, alignment=PP_ALIGN.CENTER)

add_page_number(slide5, 5)


# ═══════════════════════════════════════════
# PAGE 6 — 核心创新点
# ═══════════════════════════════════════════
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide6, BG_ROOT)
add_ambient_orbs(slide6)

add_textbox(slide6, Inches(0.9), Inches(0.45), Inches(6), Inches(0.5),
            '核心创新点', font_size=28, color=TEXT_PRI, bold=True)

add_textbox(slide6, Inches(0.9), Inches(1.05), Inches(11.5), Inches(0.45),
            '最大创新不是算法——是先证明"这条路走不通"，再找到"真正走得通的路"。',
            font_size=19, color=AMBER, bold=True)

# Three innovation cards
innov_data = [
    ('瓶颈驱动的\n技术路线', '先量化数据天花板\n（Youden\'s J / 方差分解）\n再设计多信号补偿方案\n→ 完整证据链闭环', CYAN),
    ('可解释的\n工业决策', 'SHAP + RiskDecomposer\n+ LocalExplainer\n每个告警可追溯到\n"哪个参数 → 哪个根因"', GREEN),
    ('全链路\n工程落地', '不是算法 Demo\n覆盖 数据 → 分析 → 诊断\n→ 决策 → 执行 → 反馈\n完整工业闭环', AMBER),
]
for idx, (title, desc, color) in enumerate(innov_data):
    left = Inches(0.9 + idx * 4.1)
    top = Inches(1.7)
    card = add_glass_card_colored(slide6, left, top, Inches(3.7), Inches(2.7), color)
    add_textbox(slide6, left + Inches(0.3), top + Inches(0.25), Inches(3.1), Inches(0.85),
                title, font_size=21, color=color, bold=True, line_spacing=1.3)
    add_textbox(slide6, left + Inches(0.3), top + Inches(1.2), Inches(3.1), Inches(1.3),
                desc, font_size=13, color=TEXT_SEC, line_spacing=1.55)

# Bottom cycle diagram
add_image_placeholder(slide6, Inches(2.5), Inches(4.75), Inches(8.3), Inches(2.0),
                      "Data → Insight → Decision → Action 环形流转图",
                      "传感器数据 → 统计基线 → 多信号融合 → SHAP 解释 → 工单执行 → 反馈学习")

add_page_number(slide6, 6)


# ═══════════════════════════════════════════
# PAGE 7 — 总结与展望
# ═══════════════════════════════════════════
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide7, BG_ROOT)
add_ambient_orbs(slide7)

add_textbox(slide7, Inches(0.9), Inches(0.45), Inches(6), Inches(0.5),
            '总结与展望', font_size=28, color=TEXT_PRI, bold=True)

add_textbox(slide7, Inches(0.9), Inches(1.05), Inches(11.5), Inches(0.45),
            '从100台CNC出发，构建了一套可迁移至任意工业设备的预测性维护决策范式。',
            font_size=19, color=AMBER, bold=True)

# Left: number wall (Apple KPI-style)
num_data = [
    ('100', '台 CNC', '场景规模'),
    ('9', '种故障', '全覆盖'),
    ('12', '前端页面', '全栈平台'),
    ('4', '级降级', '永不中断'),
]
for idx, (num, unit, label) in enumerate(num_data):
    row = idx // 2
    col = idx % 2
    left = Inches(0.9 + col * 2.4)
    top = Inches(1.7 + row * 1.6)
    card = add_glass_card(slide7, left, top, Inches(2.2), Inches(1.35), corner_radius=Inches(0.14))
    add_textbox(slide7, left + Inches(0.15), top + Inches(0.12), Inches(1.9), Inches(0.55),
                num, font_size=34, color=AMBER, bold=True, font_name=FONT_NUM, alignment=PP_ALIGN.CENTER)
    add_textbox(slide7, left + Inches(0.15), top + Inches(0.65), Inches(1.9), Inches(0.3),
                unit, font_size=11, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)
    add_textbox(slide7, left + Inches(0.15), top + Inches(0.95), Inches(1.9), Inches(0.25),
                label, font_size=10, color=TEXT_SEC, alignment=PP_ALIGN.CENTER, font_name=FONT_NUM)

# Middle: completed list
done_card = add_glass_card(slide7, Inches(5.9), Inches(1.7), Inches(3.5), Inches(3.0), corner_radius=Inches(0.14))
add_textbox(slide7, Inches(6.15), Inches(1.85), Inches(3.0), Inches(0.35),
            '已完成', font_size=16, color=GREEN, bold=True)
done_items = [
    '瓶颈量化：Youden\'s J / 方差分解',
    '多信号融合决策系统',
    '全栈 Web 平台（12 页面）',
    '3 策略 / 3 角色 / 4 级降级',
    'RAG 知识库 + 工单闭环',
    '故障注入 + 评委讲解助手',
]
for i, item in enumerate(done_items):
    add_textbox(slide7, Inches(6.15), Inches(2.3 + i * 0.35), Inches(3.0), Inches(0.3),
                f'·  {item}', font_size=11, color=TEXT_SEC)

# Right: future roadmap
future_card = add_glass_card(slide7, Inches(9.7), Inches(1.7), Inches(2.9), Inches(3.0), corner_radius=Inches(0.14))
add_textbox(slide7, Inches(9.95), Inches(1.85), Inches(2.4), Inches(0.35),
            '下一步', font_size=16, color=BLUE, bold=True)
future_items = [
    ('1', '多模态传感器', '振动 / 声发射 / 油液'),
    ('2', '边缘端部署', '轻量化实时推理'),
    ('3', '联邦学习', '多工厂隐私协同'),
]
for i, (step, title, desc) in enumerate(future_items):
    y = Inches(2.35 + i * 0.8)
    add_textbox(slide7, Inches(9.95), y, Inches(0.25), Inches(0.25),
                step, font_size=13, color=BLUE, bold=True, font_name=FONT_NUM)
    add_textbox(slide7, Inches(10.25), y, Inches(2.1), Inches(0.25),
                title, font_size=12, color=TEXT_PRI, bold=True)
    add_textbox(slide7, Inches(10.25), y + Inches(0.22), Inches(2.1), Inches(0.25),
                desc, font_size=10, color=TEXT_MUTED)

# Replicability banner
rep_card = add_glass_card_colored(slide7, Inches(0.9), Inches(5.1), Inches(11.7), Inches(0.65), GREEN)
add_textbox(slide7, Inches(1.2), Inches(5.2), Inches(11.1), Inches(0.45),
            '可复制性：逐设备基线方法不依赖特定传感器组合，方法论可迁移至注塑机 / 冲压机 / 风机等场景',
            font_size=14, color=TEXT_PRI, bold=True, alignment=PP_ALIGN.CENTER)

# Thank you
add_textbox(slide7, Inches(3.0), Inches(6.1), Inches(7.3), Inches(0.6),
            '谢谢各位评委', font_size=28, color=TEXT_PRI, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide7, Inches(3.0), Inches(6.7), Inches(7.3), Inches(0.35),
            '××× 团队', font_size=15, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)

add_page_number(slide7, 7)


# ═══════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════
OUTPUT = r'C:\Users\yuzzz\Desktop\苗圃杯\半决赛v1\路演PPT_预测性维护决策平台.pptx'
prs.save(OUTPUT)
print(f'PPT saved: {OUTPUT}')
print(f'Slides: {len(prs.slides)}')
